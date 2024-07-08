from pptx import Presentation
import time 
from pptx.util import Inches
 

# 生成PPT文件
def generate_ppt_file(content_list,topic,page):
    ppt=Presentation()
    
    # # PPT首页
    slide=ppt.slides.add_slide(ppt.slide_layouts[0]) # title&subtitle layout
    slide.placeholders[0].text=topic
    slide.placeholders[1].text="author:danyow"
    
    # # 内容页

    # for i,page in enumerate(ppt_content['pages']):
    for i,one_page in enumerate(content_list):
        page_list = one_page.split("\n")
        print(page_list)
        for content in page_list:
            if content[0]!='-':
                # 标题
                
                slide=ppt.slides.add_slide(ppt.slide_layouts[1]) # title&content layout
                slide.placeholders[0].text=content 
            else:
                # 内容
                sub_title=slide.placeholders[1].text_frame.add_paragraph()
                sub_title.text,sub_title.level=content[1:],1 
                # # 二级正文
                # sub_description=slide.placeholders[1].text_frame.add_paragraph()
                # sub_description.text,sub_description.level=sub_content['description'],2
    
        # break 

    file_path = f'{topic}_{int(time.time())}.pptx'
    ppt.save(file_path)
    return file_path 



def basic_usage():
    '''
    左边为文字，右边为图片'''
    ppt=Presentation()
    # first
    slide=ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.placeholders[0].text="0"
    sub_title=slide.placeholders[1].text_frame.add_paragraph() # 这个是防止出现多余的文本输入框来的
    # second 
    slide=ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.placeholders[0].text="1" 
    sub_title=slide.placeholders[1].text_frame.add_paragraph() # 这个是防止出现多余的文本输入框来的
    # 添加左侧文字占位符
    #! left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(4)) 在幻灯片上创建了一个宽 6 英寸、高 4 英寸的文本框,它距离幻灯片左边缘 0.5 英寸、顶边缘 2 英寸的位置。
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
    left_text_frame = left_box.text_frame
    left_text_frame.text = """
            1. 机器学习简介
            - 定义
            - 应用领域
            - 发展历程
            2. 机器学习基本原理
            - 监督学习
            - 无监督学习
            - 深度学习
    """
    # 添加右侧图片占位符
    #! 距左边缘 7 英寸、距顶边缘 2 英寸的图片,图片大小为宽 4 英寸、高 4 英寸。
    right_box = slide.shapes.add_picture("images/test1.png", Inches(5), Inches(2), width=Inches(4), height=Inches(4))
    
    # third 
    slide=ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.placeholders[0].text="2" 
    sub_title=slide.placeholders[1].text_frame.add_paragraph() # 这个是防止出现多余的文本输入框来的
    
    left_box = slide.shapes.add_picture("images/test1.png", Inches(0.5), Inches(2), width=Inches(4), height=Inches(4))

    
    # 添加右侧图片占位符
    #! 距左边缘 7 英寸、距顶边缘 2 英寸的图片,图片大小为宽 4 英寸、高 4 英寸。
    right_box =slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(4))
    right_box_frame = right_box.text_frame
    right_box_frame.text = """
            1. 机器学习简介
            - 定义
            - 应用领域
            - 发展历程
            2. 机器学习基本原理
            - 监督学习
            - 无监督学习
            - 深度学习
    """

    ppt.save(f'output/Test_{int(time.time())}.pptx')

    
# 生成背景自定义的PPT文件
def generate_ppt_file_bg(content_list,topic,image_path):
    ppt=Presentation()
 
    # 将图片添加到幻灯片中，并设置为背景
    left, top = Inches(0), Inches(0)
    


    # PPT首页
    slide=ppt.slides.add_slide(ppt.slide_layouts[0]) # title&subtitle layout
    slide.placeholders[0].text=topic
    slide.placeholders[1].text="author:danyow"

    # 内容页
    left = top = Inches(0)
    
    add_bg(slide,image_path,left,top,ppt)
    
    for i,one_page in enumerate(content_list):
        page_list = one_page.split("\n")        
        for content in page_list:
            if content[0]!='-':
                # 标题
                slide=ppt.slides.add_slide(ppt.slide_layouts[1]) # title&content layout
                slide.placeholders[0].text=content 
            else:
                # 内容
                sub_title=slide.placeholders[1].text_frame.add_paragraph()
                sub_title.text,sub_title.level=content[1:],1 
            


            # slide = add_bg(slide,image_path,left,top,ppt)
            add_bg(slide,image_path,left,top,ppt)


    # ppt.save('%s.pptx'%topic)
    ppt.save(f'output/Test_{topic}_{int(time.time())}.pptx')
    # ppt.save(f'output/{topic}_{int(time.time())}.pptx')


def add_bg(slide,image_path,left,top,ppt):
    pic = slide.shapes.add_picture(image_path, left, top, width=ppt.slide_width, height=ppt.slide_height)
    # This moves it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    return slide 

