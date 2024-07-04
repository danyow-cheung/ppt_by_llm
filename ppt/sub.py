from pptx import Presentation
import time 
from pptx.util import Inches,Cm 
import random 

    
class PptGenerator:
    def __init__(self) -> None:
        self.ppt=Presentation()

    def welcome_page(self,topic):
        slide=self.ppt.slides.add_slide(self.ppt.slide_layouts[0]) # title&subtitle layout
        # 首页
        slide.placeholders[0].text=topic
        slide.placeholders[1].text="author:danyow"
        return slide 
    
    def end_page(self):
        # 结束页
        slide=self.ppt.slides.add_slide(self.ppt.slide_layouts[0]) 
        slide.placeholders[0].text="Thanks for watching "
        slide.placeholders[1].text="author:danyow"
        return slide 
    def pure_text_generate(self,topic:str,content_list:list,title_list:list):

        # 首页
        slide = self.welcome_page(topic)
        # 内容页

        for i,page in enumerate(zip(content_list,title_list)):
            content = page[0]            
            title = page[1]
            slide=self.ppt.slides.add_slide(self.ppt.slide_layouts[1])

            slide.placeholders[0].text=title 
            second = slide.placeholders[1].text_frame.add_paragraph()
            second.text = content 
            second.word_wrap=True


        # 结束页
        slide = self.end_page()


        self.ppt.save(f"output/{int(time.time())}_{topic}.pptx")


        print("ppt生成完成")

    def add_bg(self,slide,image_path,ppt,left=Inches(0),top=Inches(0)):
        pic = slide.shapes.add_picture(image_path, left, top, width=ppt.slide_width, height=ppt.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        return slide 

    
    def image_text_generate(self,topic:str,content_list:list,title_list:list,images_list:list,image_list:list):
        
        # 首页
        slide = self.welcome_page(topic)
        #! 下面的情况是图片和文字长度相同
        if len(image_list)==len(content_list):
            slide = self.same_length_ppt(images_list,content_list,title_list)
        
        elif len(images_list)<len(content_list):
            raise NotImplementedError 
        
        else:
            pass 

        # 结束页
        slide = self.end_page()
        self.ppt.save(f"output/{int(time.time())}_{topic}.pptx")

    
    def same_length_ppt(self,images_list,content_list,title_list):
        for idx,value in enumerate(zip(images_list,content_list,title_list)): 
            image,text,title = value            
            if idx % 2 == 0:
                slide = self.Text2Img(text,image,title)
            else:
                # Odd
                slide = Img2Text(text,image,title)
        
        return slide 
    
    def Text2Img(self,text_content,image_path,title):
        '''
        左边为文字，右边为图片'''
        # second 
        slide=self.ppt.slides.add_slide(self.ppt.slide_layouts[1])
        slide.placeholders[0].text=title  
        sub_title=slide.placeholders[1].text_frame.add_paragraph() # 这个是防止出现多余的文本输入框来的
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
        left_text_frame = left_box.text_frame
        left_text_frame.word_wrap = True# 自动换行
        left_text_frame.text = text_content 
        image_path = rf"{image_path}"
        right_box = slide.shapes.add_picture(image_path, Inches(5), Inches(2), width=Inches(4), height=Inches(4))
        
        return slide 

    def Img2Text(self,text_content,image_path,title):
        slide=self.ppt.slides.add_slide(self.ppt.slide_layouts[1])
        slide.placeholders[0].text=title 
        sub_title=slide.placeholders[1].text_frame.add_paragraph() # 这个是防止出现多余的文本输入框来的
        image_path = rf"{image_path}"
        left_box = slide.shapes.add_picture(image_path, Inches(0.5), Inches(2), width=Inches(4), height=Inches(4))
        # 添加右侧图片占位符
        #! 距左边缘 7 英寸、距顶边缘 2 英寸的图片,图片大小为宽 4 英寸、高 4 英寸。
        right_box =slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(4))
        right_box_frame = right_box.text_frame
        right_box_frame.word_wrap = True # 自动换行
        right_box_frame.text = text_content
        return slide 
